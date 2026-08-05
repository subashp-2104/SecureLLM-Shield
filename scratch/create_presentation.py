import sys
import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette
    COLOR_BG = RGBColor(11, 19, 43)        # #0B132B Deep Navy
    COLOR_CARD = RGBColor(28, 37, 65)      # #1C2541 Dark Slate Blue
    COLOR_BORDER = RGBColor(0, 242, 254)   # #00F2FE Electric Cyan
    COLOR_CYAN = RGBColor(0, 242, 254)     # #00F2FE
    COLOR_TEAL = RGBColor(72, 202, 228)    # #48CAE4
    COLOR_WHITE = RGBColor(255, 255, 255)  # White
    COLOR_MUTED = RGBColor(141, 153, 174)  # Slate Gray
    COLOR_GREEN = RGBColor(57, 255, 20)    # Neon Green
    COLOR_RED = RGBColor(255, 0, 85)       # Neon Red
    COLOR_GOLD = RGBColor(255, 215, 0)     # Gold

    # Image Paths
    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    IMG_SHIELD = os.path.join(base_dir, "assets", "security_shield.jpg")
    IMG_NODES = os.path.join(base_dir, "assets", "threat_nodes.jpg")
    IMG_SAEC_LOGO = os.path.join(base_dir, "assets", "saec_logo.png")

    def add_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, category_text="SECURELLM SHIELD | PROJECT REVIEW"):
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_CYAN
        p_cat.font.name = "Calibri"

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_WHITE
        p_title.font.name = "Calibri"

    def add_footer(slide, current_slide, total_slides=19):
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.1), Inches(11.7), Inches(0.35))
        tf = footer_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"SecureLLM Shield Presentation  |  Department of CSE (AI & ML)  |  Slide {current_slide} of {total_slides}"
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_MUTED
        p.font.name = "Calibri"

    def set_speaker_notes(slide, notes_text):
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = notes_text

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide (With SAEC Logo Banner & Security Shield)
    # -------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    add_bg(s1)

    # Top College Logo Banner Container Box (White Fill)
    logo_banner = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.4), Inches(11.733), Inches(1.15))
    logo_banner.fill.solid()
    logo_banner.fill.fore_color.rgb = COLOR_WHITE
    logo_banner.line.color.rgb = COLOR_CYAN
    logo_banner.line.width = Pt(1.5)

    if os.path.exists(IMG_SAEC_LOGO):
        s1.shapes.add_picture(IMG_SAEC_LOGO, Inches(0.95), Inches(0.45), Inches(11.433), Inches(1.05))

    # Main Project Presentation Container Card
    card1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.7), Inches(11.733), Inches(5.3))
    card1.fill.solid()
    card1.fill.fore_color.rgb = COLOR_CARD
    card1.line.color.rgb = COLOR_CYAN
    card1.line.width = Pt(2)

    tb = s1.shapes.add_textbox(Inches(1.1), Inches(1.85), Inches(7.3), Inches(4.9))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "SecureLLM Shield"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    p2 = tf.add_paragraph()
    p2.text = "Privacy Threat Detection, Sensitive Information Protection, Audit and Migration Framework for Large Language Models"
    p2.font.size = Pt(13)
    p2.font.color.rgb = COLOR_WHITE
    p2.space_before = Pt(4)

    p3 = tf.add_paragraph()
    p3.text = "DEPARTMENT OF CSE (ARTIFICIAL INTELLIGENCE & MACHINE LEARNING)"
    p3.font.size = Pt(11)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_TEAL
    p3.space_before = Pt(10)

    # Project Guide Section
    p_g1 = tf.add_paragraph()
    p_g1.text = "Project Guide:"
    p_g1.font.size = Pt(11.5)
    p_g1.font.bold = True
    p_g1.font.color.rgb = COLOR_GOLD
    p_g1.space_before = Pt(10)

    p_g2 = tf.add_paragraph()
    p_g2.text = "Mr. G. Saravanan Gokul, M.E.\nAssistant Professor, Dept. of CSE (AI & ML), S.A. Engineering College"
    p_g2.font.size = Pt(12)
    p_g2.font.bold = True
    p_g2.font.color.rgb = COLOR_WHITE
    p_g2.space_before = Pt(2)

    # Team Members Section
    p_t1 = tf.add_paragraph()
    p_t1.text = "Presented By:"
    p_t1.font.size = Pt(11.5)
    p_t1.font.bold = True
    p_t1.font.color.rgb = COLOR_CYAN
    p_t1.space_before = Pt(10)

    p_t2 = tf.add_paragraph()
    p_t2.text = "BARATHKUMAR M (Reg. No: 111923AM01007)  |  SUBASH P (Reg. No: 111923AM02052)"
    p_t2.font.size = Pt(12)
    p_t2.font.bold = True
    p_t2.font.color.rgb = COLOR_WHITE
    p_t2.space_before = Pt(2)

    p_t3 = tf.add_paragraph()
    p_t3.text = "Academic Project Review — 2026"
    p_t3.font.size = Pt(10)
    p_t3.font.color.rgb = COLOR_MUTED
    p_t3.space_before = Pt(6)

    if os.path.exists(IMG_SHIELD):
        s1.shapes.add_picture(IMG_SHIELD, Inches(8.6), Inches(1.85), Inches(3.7), Inches(4.95))

    set_speaker_notes(s1, "Good morning respected evaluators. Today we are presenting our project titled SecureLLM Shield under the guidance of Mr. G. Saravanan Gokul, Assistant Professor, Department of CSE (AI & ML), S.A. Engineering College.")

    # -------------------------------------------------------------
    # SLIDE 2: Team Members & Project Guide Credentials
    # -------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    add_bg(s2)
    add_header(s2, "Team Members & Project Guide Credentials")
    add_footer(s2, 2)

    # Card 1: BARATHKUMAR M
    c1 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.7), Inches(3.65), Inches(4.8))
    c1.fill.solid()
    c1.fill.fore_color.rgb = COLOR_CARD
    c1.line.color.rgb = COLOR_CYAN
    c1.line.width = Pt(1.5)

    tb1 = s2.shapes.add_textbox(Inches(0.95), Inches(1.85), Inches(3.35), Inches(4.5))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "BARATHKUMAR M"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    items1 = [
        "Reg. No: 111923AM01007",
        "Dept: CSE (AI & ML)",
        "Role: Security Pipeline Developer",
        "Responsibilities:",
        "  • 20-Entity PII Pattern Engines",
        "  • Aadhaar, PAN & Bank Masking",
        "  • Jailbreak Defense Guardrails"
    ]
    for item in items1:
        pi = tf1.add_paragraph()
        pi.text = item
        pi.font.size = Pt(11.5)
        pi.font.bold = ("Reg. No" in item or "Role:" in item or "Responsibilities" in item)
        pi.font.color.rgb = COLOR_WHITE if "  •" in item else (COLOR_TEAL if "Role:" in item or "Responsibilities:" in item else COLOR_MUTED)
        pi.space_before = Pt(5)

    # Card 2: SUBASH P
    c2 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.84), Inches(1.7), Inches(3.65), Inches(4.8))
    c2.fill.solid()
    c2.fill.fore_color.rgb = COLOR_CARD
    c2.line.color.rgb = COLOR_CYAN
    c2.line.width = Pt(1.5)

    tb2 = s2.shapes.add_textbox(Inches(4.99), Inches(1.85), Inches(3.35), Inches(4.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True

    p = tf2.paragraphs[0]
    p.text = "SUBASH P"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    items2 = [
        "Reg. No: 111923AM02052",
        "Dept: CSE (AI & ML)",
        "Role: Backend & UI Developer",
        "Responsibilities:",
        "  • Python Flask API Architecture",
        "  • Interactive Security Dashboard",
        "  • SHA-256 Audit Ledger & Copilot"
    ]
    for item in items2:
        pi = tf2.add_paragraph()
        pi.text = item
        pi.font.size = Pt(11.5)
        pi.font.bold = ("Reg. No" in item or "Role:" in item or "Responsibilities" in item)
        pi.font.color.rgb = COLOR_WHITE if "  •" in item else (COLOR_TEAL if "Role:" in item or "Responsibilities:" in item else COLOR_MUTED)
        pi.space_before = Pt(5)

    # Card 3: PROJECT GUIDE
    c3 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.88), Inches(1.7), Inches(3.65), Inches(4.8))
    c3.fill.solid()
    c3.fill.fore_color.rgb = COLOR_CARD
    c3.line.color.rgb = COLOR_GOLD
    c3.line.width = Pt(1.5)

    tb3 = s2.shapes.add_textbox(Inches(9.03), Inches(1.85), Inches(3.35), Inches(4.5))
    tf3 = tb3.text_frame
    tf3.word_wrap = True

    p = tf3.paragraphs[0]
    p.text = "PROJECT GUIDE"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD

    items3 = [
        "Mr. G. Saravanan Gokul, M.E.",
        "Desig: Assistant Professor",
        "Dept: CSE (AI & ML)",
        "Institution: S.A. Engineering College",
        "Guidance Scope:",
        "  • Security Architecture Oversight",
        "  • PII & Threat Defense Review",
        "  • Technical Research Guidance"
    ]
    for item in items3:
        pi = tf3.add_paragraph()
        pi.text = item
        pi.font.size = Pt(11.5)
        pi.font.bold = ("Mr. G." in item or "Desig:" in item or "Guidance Scope:" in item)
        pi.font.color.rgb = COLOR_WHITE if "  •" in item else (COLOR_TEAL if "Desig:" in item or "Guidance Scope:" in item else COLOR_MUTED)
        pi.space_before = Pt(5)

    set_speaker_notes(s2, "This presentation is prepared by Barathkumar M and Subash P. Barathkumar handled the security pipeline and PII detection algorithms, while Subash built the Python backend and the interactive dashboard.")

    # -------------------------------------------------------------
    # SLIDE 3: Project Overview (With Graphic Image)
    # -------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    add_bg(s3)
    add_header(s3, "Project Overview")
    add_footer(s3, 3)

    points3 = [
        ("AI Security Gateway:", "Acts as an inline security barrier placed between enterprise users and external LLM models."),
        ("Universal PII Masking:", "Scans prompts for 20 sensitive entity types (Aadhaar, PAN, Bank Accounts) and applies partial masking."),
        ("Prompt Injection Defense:", "Detects and blocks jailbreak attacks (DAN prompts) before they reach foundation models."),
        ("Risk Scoring Engine:", "Calculates real-time risk scores (0% to 100%) and displays explainable AI (XAI) security decision logs."),
        ("Cryptographic Audit Ledger:", "Stores security events in a tamper-proof SHA-256 blockchain audit trail for compliance.")
    ]

    for idx, (title, desc) in enumerate(points3):
        box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6 + idx * 1.02), Inches(7.8), Inches(0.88))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD
        box.line.color.rgb = COLOR_CYAN
        box.line.width = Pt(1)

        tb = s3.shapes.add_textbox(Inches(1.0), Inches(1.68 + idx * 1.02), Inches(7.4), Inches(0.72))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = title + " "
        r1.font.bold = True
        r1.font.size = Pt(14)
        r1.font.color.rgb = COLOR_CYAN

        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(13)
        r2.font.color.rgb = COLOR_WHITE

    if os.path.exists(IMG_SHIELD):
        s3.shapes.add_picture(IMG_SHIELD, Inches(8.9), Inches(1.6), Inches(3.6), Inches(5.0))

    set_speaker_notes(s3, "In simple terms, SecureLLM Shield acts as a protective shield between users and AI models like ChatGPT. It ensures no secret or personal data leaks out while asking questions to AI.")

    # -------------------------------------------------------------
    # SLIDE 4: Problem Statement
    # -------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    add_bg(s4)
    add_header(s4, "Problem Statement")
    add_footer(s4, 4)

    prob_points = [
        "Unintentional Data Leakage: Employees frequently paste secret source code, customer details, and identity numbers into AI tools.",
        "Privacy Violations: Public LLM servers store prompt history, risking exposure of confidential citizen data.",
        "Failure with Indian Identity Tags: Standard regex filters fail on Indian 12-digit Aadhaar cards, PAN numbers, and IFSC codes.",
        "Adversarial Jailbreak Attacks: Hackers use prompt injection tricks (DAN payloads) to bypass safety guardrails.",
        "Lack of Audit Trail: Organizations have no centralized system to track AI usage, calculate risk, or prove regulatory compliance."
    ]

    for idx, prob in enumerate(prob_points):
        box = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6 + idx * 1.02), Inches(11.7), Inches(0.88))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD
        box.line.color.rgb = COLOR_RED
        box.line.width = Pt(1)

        tb = s4.shapes.add_textbox(Inches(1.1), Inches(1.68 + idx * 1.02), Inches(11.1), Inches(0.72))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        parts = prob.split(":", 1)
        
        r1 = p.add_run()
        r1.text = "⚠️ " + parts[0] + ":"
        r1.font.bold = True
        r1.font.size = Pt(15)
        r1.font.color.rgb = COLOR_RED

        r2 = p.add_run()
        r2.text = parts[1]
        r2.font.size = Pt(14)
        r2.font.color.rgb = COLOR_WHITE

    set_speaker_notes(s4, "The main problem today is that users unknowingly share personal data like Aadhaar, passwords, or financial details with AI tools. Existing filters fail to catch these Indian identity formats, creating high security risks.")

    # -------------------------------------------------------------
    # SLIDE 5: Project Objectives
    # -------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    add_bg(s5)
    add_header(s5, "Project Objectives")
    add_footer(s5, 5)

    objs = [
        "1. Build a Real-Time Prompt Sanitizer: Detect and sanitize 20 sensitive entity types in under 50 milliseconds.",
        "2. Implement Context-Preserving Partial Masking: Hide critical digits (e.g., XXXX XXXX 9012) while keeping prompt meaning intact.",
        "3. Neutralize Adversarial Prompt Injections: Block DAN jailbreaks, system overrides, and context poisoning attempts.",
        "4. Develop AI-Based Risk Prediction: Assign real-time threat scores (0% to 100%) with explainable decision logs.",
        "5. Guarantee Immutable Auditability: Store tamper-proof security logs using a cryptographic SHA-256 blockchain ledger."
    ]

    for idx, obj in enumerate(objs):
        box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6 + idx * 1.02), Inches(11.7), Inches(0.88))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD
        box.line.color.rgb = COLOR_CYAN
        box.line.width = Pt(1)

        tb = s5.shapes.add_textbox(Inches(1.1), Inches(1.68 + idx * 1.02), Inches(11.1), Inches(0.72))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        parts = obj.split(":", 1)

        r1 = p.add_run()
        r1.text = parts[0] + ":"
        r1.font.bold = True
        r1.font.size = Pt(15)
        r1.font.color.rgb = COLOR_CYAN

        r2 = p.add_run()
        r2.text = parts[1]
        r2.font.size = Pt(14)
        r2.font.color.rgb = COLOR_WHITE

    set_speaker_notes(s5, "Our objective is to build a smart sanitizer that hides personal details while keeping the text understandable for AI. We also calculate real-time risk scores and block hacker attacks.")

    # -------------------------------------------------------------
    # SLIDE 6: Existing System vs Limitations
    # -------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    add_bg(s6)
    add_header(s6, "Existing System vs Limitations")
    add_footer(s6, 6)

    c_left = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.7), Inches(5.6), Inches(4.8))
    c_left.fill.solid()
    c_left.fill.fore_color.rgb = COLOR_CARD
    c_left.line.color.rgb = COLOR_MUTED
    c_left.line.width = Pt(1)

    tb_l = s6.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(5.2), Inches(4.4))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    p = tf_l.paragraphs[0]
    p.text = "Existing Security Tools"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    l_items = [
        "• Basic Keyword Matching (Simple Regex)",
        "• Full Redaction ([REDACTED] or [BLOCKED])",
        "• Static Single-Provider Rules",
        "• Manual Security Inspection",
        "• Standard Log Files (Editable & Unsecured)"
    ]
    for item in l_items:
        pi = tf_l.add_paragraph()
        pi.text = item
        pi.font.size = Pt(14)
        pi.font.color.rgb = COLOR_WHITE
        pi.space_before = Pt(12)

    c_right = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.7), Inches(5.7), Inches(4.8))
    c_right.fill.solid()
    c_right.fill.fore_color.rgb = COLOR_CARD
    c_right.line.color.rgb = COLOR_RED
    c_right.line.width = Pt(1.5)

    tb_r = s6.shapes.add_textbox(Inches(7.0), Inches(1.9), Inches(5.3), Inches(4.4))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    p = tf_r.paragraphs[0]
    p.text = "Major Limitations"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED

    r_items = [
        "❌ Misses Formatted Indian PII (Aadhaar / PAN)",
        "❌ Destroys Sentence Context for AI Models",
        "❌ Vulnerable to DAN Prompt Injection Jailbreaks",
        "❌ High Processing Delays (> 500ms)",
        "❌ No Immutable Auditability for Compliance"
    ]
    for item in r_items:
        pi = tf_r.add_paragraph()
        pi.text = item
        pi.font.size = Pt(14)
        pi.font.color.rgb = COLOR_WHITE
        pi.space_before = Pt(12)

    set_speaker_notes(s6, "Current security systems either block everything using basic keyword matching or erase full words, which makes the AI answer useless. They also don't recognize Indian identity formats.")

    # -------------------------------------------------------------
    # SLIDE 7: Comparison with Existing IEEE Conference Research (NEW)
    # -------------------------------------------------------------
    s7_ieee = prs.slides.add_slide(blank_layout)
    add_bg(s7_ieee)
    add_header(s7_ieee, "Comparison with Existing IEEE Conference Research", "SECURELLM SHIELD | IEEE RESEARCH BENCHMARKING")
    add_footer(s7_ieee, 7)

    # Left Box: Existing IEEE Research Limitations
    box_ieee_l = s7_ieee.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.38), Inches(5.7), Inches(3.3))
    box_ieee_l.fill.solid()
    box_ieee_l.fill.fore_color.rgb = COLOR_CARD
    box_ieee_l.line.color.rgb = COLOR_RED
    box_ieee_l.line.width = Pt(1.5)

    tb_ieee_l = s7_ieee.shapes.add_textbox(Inches(1.0), Inches(1.48), Inches(5.3), Inches(3.1))
    tf_ieee_l = tb_ieee_l.text_frame
    tf_ieee_l.word_wrap = True
    p = tf_ieee_l.paragraphs[0]
    p.text = "Existing IEEE Research Limitations"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED

    ieee_limitations = [
        "• Focuses only on basic PII regex matching",
        "• Supports text inputs with no jailbreak defense",
        "• No response safety verification or output validation",
        "• Lacks security middleware proxy architecture",
        "• No protection for Indian Aadhaar, PAN or API Keys",
        "• Proof-of-concept focus with no immutable audit logging"
    ]
    for item in ieee_limitations:
        pi = tf_ieee_l.add_paragraph()
        pi.text = item
        pi.font.size = Pt(10.5)
        pi.font.color.rgb = COLOR_WHITE
        pi.space_before = Pt(4)

    # Right Box: How SecureLLM Shield Improves These Limitations
    box_ieee_r = s7_ieee.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.833), Inches(1.38), Inches(5.7), Inches(3.3))
    box_ieee_r.fill.solid()
    box_ieee_r.fill.fore_color.rgb = COLOR_CARD
    box_ieee_r.line.color.rgb = COLOR_CYAN
    box_ieee_r.line.width = Pt(1.5)

    tb_ieee_r = s7_ieee.shapes.add_textbox(Inches(7.033), Inches(1.48), Inches(5.3), Inches(3.1))
    tf_ieee_r = tb_ieee_r.text_frame
    tf_ieee_r.word_wrap = True
    p = tf_ieee_r.paragraphs[0]
    p.text = "How SecureLLM Shield Improves Limitations"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    shield_improvements = [
        "✅ Multi-layer security architecture (Regex + NER Spans)",
        "✅ Automatic partial masking of confidential data & secrets",
        "✅ Prompt injection detection & DAN jailbreak filtering",
        "✅ Pre-LLM sanitization + post-LLM output verification",
        "✅ Security middleware proxy between User and LLM",
        "✅ Comprehensive Aadhaar, PAN, Bank, Password & API Key protection",
        "✅ Scalable architecture compatible with GPT, Llama, Gemini, Mistral"
    ]
    for item in shield_improvements:
        pi = tf_ieee_r.add_paragraph()
        pi.text = item
        pi.font.size = Pt(10.5)
        pi.font.color.rgb = COLOR_WHITE
        pi.space_before = Pt(3.5)

    # Bottom Comparison Table
    t_shape_ieee = s7_ieee.shapes.add_table(3, 4, Inches(0.8), Inches(4.78), Inches(11.733), Inches(1.2))
    table_ieee = t_shape_ieee.table
    table_ieee.columns[0].width = Inches(2.933)
    table_ieee.columns[1].width = Inches(2.933)
    table_ieee.columns[2].width = Inches(2.933)
    table_ieee.columns[3].width = Inches(2.934)

    t_matrix = [
        [("Existing Research", "Basic PII Detection"), ("SecureLLM Shield", "Multi-layer Privacy Protection"), ("Existing Research", "Limited Security Scope"), ("SecureLLM Shield", "End-to-End Security Pipeline")],
        [("Existing Research", "No Output Validation"), ("SecureLLM Shield", "Output Safety Verification"), ("Existing Research", "Basic Prompt Filtering"), ("SecureLLM Shield", "Advanced Prompt Validation")],
        [("Existing Research", "Direct LLM Access"), ("SecureLLM Shield", "Secure Middleware Layer"), ("Existing Research", "Proof-of-Concept System"), ("SecureLLM Shield", "Privacy-First Architecture + Audit")]
    ]

    for r_idx, row_data in enumerate(t_matrix):
        for c_idx, (sys_tag, sys_val) in enumerate(row_data):
            cell = table_ieee.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(18, 25, 48) if "Existing" in sys_tag else COLOR_CARD
            p = cell.text_frame.paragraphs[0]
            
            r1 = p.add_run()
            r1.text = sys_tag + ": "
            r1.font.bold = True
            r1.font.size = Pt(9.5)
            r1.font.color.rgb = COLOR_RED if "Existing" in sys_tag else COLOR_CYAN

            r2 = p.add_run()
            r2.text = sys_val
            r2.font.size = Pt(10)
            r2.font.color.rgb = COLOR_WHITE

    # Conclusion Box at Bottom
    conc_box_ieee = s7_ieee.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.10), Inches(11.733), Inches(0.85))
    conc_box_ieee.fill.solid()
    conc_box_ieee.fill.fore_color.rgb = COLOR_CARD
    conc_box_ieee.line.color.rgb = COLOR_CYAN
    conc_box_ieee.line.width = Pt(1.5)

    tb_conc = s7_ieee.shapes.add_textbox(Inches(0.95), Inches(6.14), Inches(11.433), Inches(0.77))
    tf_conc = tb_conc.text_frame
    tf_conc.word_wrap = True
    p = tf_conc.paragraphs[0]
    p.text = "📌 Note: Our Secure LLM Shield is not intended to replace previous IEEE research. Instead, it extends existing approaches by integrating multiple security mechanisms into a single end-to-end framework for safer LLM interactions."
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEAL

    set_speaker_notes(s7_ieee, "When comparing our project to existing IEEE conference research, prior works typically focus on basic regex PII filters or static prompt rules. SecureLLM Shield extends these approaches by introducing an end-to-end security middleware with dual scanning, jailbreak defense, output verification, and cryptographic auditability.")

    # -------------------------------------------------------------
    # SLIDE 8: Proposed System (Shifted to Slide 8)
    # -------------------------------------------------------------
    s8 = prs.slides.add_slide(blank_layout)
    add_bg(s8)
    add_header(s8, "Proposed System — SecureLLM Shield")
    add_footer(s8, 8)

    prop_cards = [
        ("Multi-Stage Pipeline", "Regex + NER + Classifier", "Scans text through high-speed pattern matching and AI entity spans."),
        ("Smart Partial Masking", "Context Preserved", "Replaces sensitive digits while preserving format (e.g. XXXX XXXX 9012)."),
        ("Indian PII Coverage", "Universal Protection", "Full support for Aadhaar, PAN, Bank Accounts, UPI, and IFSC codes."),
        ("Adaptive Policy Engine", "Industry Presets", "Configurable profiles for Enterprise, Healthcare, Banking, and Government."),
        ("Cryptographic Audit", "SHA-256 Blockchain", "Stores tamper-proof security audit logs in an immutable ledger."),
        ("Security Copilot", "AI Assistant Drawer", "Embedded helper that explains security risks and compliance rules live.")
    ]

    for idx, (title, tag, desc) in enumerate(prop_cards):
        row = idx // 3
        col = idx % 3
        x = Inches(0.8 + col * 3.98)
        y = Inches(1.7 + row * 2.5)

        box = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.75), Inches(2.2))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD
        box.line.color.rgb = COLOR_CYAN
        box.line.width = Pt(1.5)

        tb = s8.shapes.add_textbox(x + Inches(0.15), y + Inches(0.15), Inches(3.45), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_CYAN

        p2 = tf.add_paragraph()
        p2.text = tag
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_TEAL
        p2.space_before = Pt(2)

        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(12)
        p3.font.color.rgb = COLOR_WHITE
        p3.space_before = Pt(6)

    set_speaker_notes(s8, "Our proposed system solves these issues by using a hybrid pipeline. It masks only the secret parts of numbers—like showing only the last 4 digits of Aadhaar—so the AI still understands the context.")

    # -------------------------------------------------------------
    # SLIDE 9: System Architecture (With Diagram Image)
    # -------------------------------------------------------------
    s9 = prs.slides.add_slide(blank_layout)
    add_bg(s9)
    add_header(s9, "System Architecture")
    add_footer(s9, 9)

    box_data = [
        ("1. User Input", "Prompt containing text, code, or identity tags"),
        ("2. Security Gateway", "SecureLLM Shield inline inspection proxy"),
        ("3. Hybrid Pipeline", "Regex + NER + Injection Classifier"),
        ("4. Mask & Score", "Partial Masking + Risk Score calculation"),
        ("5. Target LLM", "Sanitized prompt sent to GPT-4o / Claude / Llama")
    ]

    for idx, (b_title, b_desc) in enumerate(box_data):
        x = Inches(0.8 + idx * 2.4)
        y = Inches(1.6)

        box = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.1), Inches(2.2))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD
        box.line.color.rgb = COLOR_CYAN
        box.line.width = Pt(1.5)

        tb = s9.shapes.add_textbox(x + Inches(0.1), y + Inches(0.12), Inches(1.9), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = b_title
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_CYAN
        p1.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = b_desc
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = COLOR_WHITE
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(8)

        if idx < 4:
            arrow = s9.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(2.12), Inches(2.5), Inches(0.24), Inches(0.35))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLOR_CYAN
            arrow.line.fill.background()

    if os.path.exists(IMG_SHIELD):
        s9.shapes.add_picture(IMG_SHIELD, Inches(0.8), Inches(4.1), Inches(11.7), Inches(2.6))

    set_speaker_notes(s9, "Here is our System Architecture. The user prompt goes first into our gateway. The gateway scans, masks sensitive data, checks for attacks, and only then forwards the safe prompt to the AI model.")

    # -------------------------------------------------------------
    # SLIDE 10: Modules of the Project
    # -------------------------------------------------------------
    s10 = prs.slides.add_slide(blank_layout)
    add_bg(s10)
    add_header(s10, "Modules of the Project")
    add_footer(s10, 10)

    modules = [
        ("Module 1", "Universal PII Detection Engine", "Scans text for 20 sensitive entity types including Indian Aadhaar, PAN, Bank Accounts, Passports, and API Keys."),
        ("Module 2", "Prompt Injection & Jailbreak Defense", "Detects Do-Anything-Now (DAN) jailbreaks, system prompt overrides, and context poisoning attempts."),
        ("Module 3", "Adaptive Privacy Policy Engine", "Allows administrators to switch preset profiles (Enterprise, Healthcare, Banking, Government) or toggle custom rules."),
        ("Module 4", "Risk Scoring & XAI Decision Dashboard", "Calculates live risk scores (0% to 100%) and displays explainable AI decision logs with highlighted diffs."),
        ("Module 5", "Cryptographic Audit & Migration Framework", "Records security events in an immutable SHA-256 blockchain ledger and clones policies across LLM providers.")
    ]

    for idx, (m_id, m_name, m_desc) in enumerate(modules):
        box = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6 + idx * 1.02), Inches(11.7), Inches(0.88))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD
        box.line.color.rgb = COLOR_CYAN
        box.line.width = Pt(1)

        tb = s10.shapes.add_textbox(Inches(1.1), Inches(1.68 + idx * 1.02), Inches(11.1), Inches(0.72))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        
        r1 = p.add_run()
        r1.text = f"[{m_id}] {m_name}: "
        r1.font.bold = True
        r1.font.size = Pt(15)
        r1.font.color.rgb = COLOR_CYAN

        r2 = p.add_run()
        r2.text = m_desc
        r2.font.size = Pt(13)
        r2.font.color.rgb = COLOR_WHITE

    set_speaker_notes(s10, "Our project is divided into 5 core modules: PII Detection, Jailbreak Defense, Policy Configuration, Risk Scoring, and Cryptographic Logging.")

    # -------------------------------------------------------------
    # SLIDE 11: PII Detection & Partial Masking
    # -------------------------------------------------------------
    s11 = prs.slides.add_slide(blank_layout)
    add_bg(s11)
    add_header(s11, "PII Detection & Partial Masking")
    add_footer(s11, 11)

    pii_examples = [
        ("Aadhaar Number", "4567 8912 3456", "XXXX XXXX 3456", "Preserves last 4 digits for audit"),
        ("PAN Card ID", "ABCDE1234F", "XXXXX1234F", "Conceals identity letter prefix"),
        ("Bank Account", "12345678901", "XXXXXXX8901", "PCI-DSS compliant 4-digit mask"),
        ("UPI Payment ID", "sonuz@oksbi", "s****@oksbi", "Preserves initial & payment domain"),
        ("Developer API Key", "sk_live_51Nz84...", "[API_KEY_REDACTED]", "Full secret token redaction")
    ]

    headers = ["Entity Type", "Original Input", "Sanitized Output", "Masking Strategy"]
    widths = [Inches(2.5), Inches(2.8), Inches(2.8), Inches(3.4)]

    t_shape = s11.shapes.add_table(6, 4, Inches(0.8), Inches(1.7), Inches(11.7), Inches(4.8))
    table = t_shape.table

    for i, w in enumerate(widths):
        table.columns[i].width = w

    for c_idx, h_text in enumerate(headers):
        cell = table.cell(0, c_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_CARD
        p = cell.text_frame.paragraphs[0]
        p.text = h_text
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_CYAN

    for r_idx, row_data in enumerate(pii_examples):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(18, 25, 48) if r_idx % 2 == 0 else COLOR_CARD
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(13)
            p.font.color.rgb = COLOR_GREEN if c_idx == 2 else COLOR_WHITE

    set_speaker_notes(s11, "This slide demonstrates our partial masking feature. Notice how Aadhaar and PAN numbers keep their format so the AI understands it is an ID, without revealing the actual private number.")

    # -------------------------------------------------------------
    # SLIDE 12: Prompt Injection & Defense (With Threat Nodes Graphic)
    # -------------------------------------------------------------
    s12 = prs.slides.add_slide(blank_layout)
    add_bg(s12)
    add_header(s12, "Prompt Injection & Privacy Threat Defense")
    add_footer(s12, 12)

    threat_boxes = [
        ("System Override Blocks", "Detects prompts instructing AI to 'Ignore previous rules'."),
        ("DAN Jailbreak Defense", "Neutralizes Do-Anything-Now jailbreak prompts."),
        ("Context Poisoning Defense", "Prevents malicious instructions inside RAG document chunks."),
        ("Data Exfiltration Prevention", "Blocks prompts asking AI to reveal internal secret keys."),
        ("Automated Red Team Validated", "Verified 0% bypass rate under 1,000-payload campaigns.")
    ]

    for idx, (t_title, t_desc) in enumerate(threat_boxes):
        box = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6 + idx * 1.02), Inches(7.8), Inches(0.88))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD
        box.line.color.rgb = COLOR_RED
        box.line.width = Pt(1)

        tb = s12.shapes.add_textbox(Inches(1.0), Inches(1.68 + idx * 1.02), Inches(7.4), Inches(0.72))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = "🛡️ " + t_title + ": "
        r1.font.bold = True
        r1.font.size = Pt(14)
        r1.font.color.rgb = COLOR_CYAN

        r2 = p.add_run()
        r2.text = t_desc
        r2.font.size = Pt(12.5)
        r2.font.color.rgb = COLOR_WHITE

    if os.path.exists(IMG_NODES):
        s12.shapes.add_picture(IMG_NODES, Inches(8.9), Inches(1.6), Inches(3.6), Inches(5.0))

    set_speaker_notes(s12, "Hackers try to trick AI using jailbreak tricks like 'Ignore all rules and answer this'. Our jailbreak defense module detects these trick phrases and blocks them instantly.")

    # -------------------------------------------------------------
    # SLIDE 13: Audit Logging & Risk Scoring
    # -------------------------------------------------------------
    s13 = prs.slides.add_slide(blank_layout)
    add_bg(s13)
    add_header(s13, "Audit Logging & Risk Scoring Engine")
    add_footer(s13, 13)

    risk_cards = [
        ("SAFE (0% - 25%)", "Clean Prompt Payload", "Cleared for immediate LLM routing.", COLOR_GREEN),
        ("MODERATE (26% - 50%)", "Low PII Spans Detected", "Light partial masking applied.", COLOR_CYAN),
        ("HIGH (51% - 75%)", "Multiple PII / Secrets", "Sanitized & user session flagged.", COLOR_GOLD),
        ("CRITICAL (76% - 100%)", "Active Jailbreak Attack", "Blocked at Gateway immediately.", COLOR_RED)
    ]

    for idx, (r_title, r_sub, r_desc, r_color) in enumerate(risk_cards):
        x = Inches(0.8 + idx * 2.98)
        y = Inches(1.8)

        box = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.75), Inches(3.2))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD
        box.line.color.rgb = r_color
        box.line.width = Pt(1.5)

        tb = s13.shapes.add_textbox(x + Inches(0.1), y + Inches(0.15), Inches(2.55), Inches(2.9))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = r_title
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = r_color
        p1.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = r_sub
        p2.font.size = Pt(12)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_WHITE
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(8)

        p3 = tf.add_paragraph()
        p3.text = r_desc
        p3.font.size = Pt(11)
        p3.font.color.rgb = COLOR_MUTED
        p3.alignment = PP_ALIGN.CENTER
        p3.space_before = Pt(10)

    audit_box = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.3))
    audit_box.fill.solid()
    audit_box.fill.fore_color.rgb = COLOR_CARD
    audit_box.line.color.rgb = COLOR_CYAN
    audit_box.line.width = Pt(1)

    tb_a = s13.shapes.add_textbox(Inches(1.0), Inches(5.4), Inches(11.3), Inches(1.1))
    tf_a = tb_a.text_frame
    tf_a.word_wrap = True

    p = tf_a.paragraphs[0]
    p.text = "Cryptographic Blockchain Audit Ledger:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    p_a2 = tf_a.add_paragraph()
    p_a2.text = "• Every request evaluation creates an immutable block linked with SHA-256 cryptographic hashes.\n• Allows security teams to click any status badge (Masked / Blocked / Cleared) for full audit trace inspection."
    p_a2.font.size = Pt(13)
    p_a2.font.color.rgb = COLOR_WHITE
    p_a2.space_before = Pt(4)

    set_speaker_notes(s13, "Every prompt receives a live Risk Score between 0 and 100 percent based on entity count and threat level. Critical threats are blocked automatically and recorded in our tamper-proof log.")

    # -------------------------------------------------------------
    # SLIDE 14: Technologies Used
    # -------------------------------------------------------------
    s14 = prs.slides.add_slide(blank_layout)
    add_bg(s14)
    add_header(s14, "Technologies Used")
    add_footer(s14, 14)

    techs = [
        ("Backend Framework", "Python 3.12 & Flask Server", "RESTful API handling prompt sanitizer routes and health endpoints."),
        ("Frontend & Dashboard", "HTML5, Vanilla CSS3, JavaScript", "Sleek dark-mode interface with multi-color glowing themes (Cyber, Vibrant, Slate)."),
        ("Detection Engines", "Regex Pattern Matching + NER Spans", "Universal 20-entity detection rules including Aadhaar, PAN, and Bank details."),
        ("Security & VectorDB", "FAISS VectorDB + SHA-256 Hashes", "Role-Based Access Control (RBAC) and immutable cryptographic audit logging."),
        ("Data Visualization", "Chart.js Library", "Real-time 24h threat timeline and model safety vs latency bubble charts."),
        ("Development Tools", "VS Code & Git Repository", "Hosted at github.com/subashp-2104/SecureLLM-Shield for team collaboration.")
    ]

    for idx, (cat, title, desc) in enumerate(techs):
        row = idx // 3
        col = idx % 3
        x = Inches(0.8 + col * 3.98)
        y = Inches(1.7 + row * 2.5)

        box = s14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.75), Inches(2.2))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD
        box.line.color.rgb = COLOR_CYAN
        box.line.width = Pt(1.5)

        tb = s14.shapes.add_textbox(x + Inches(0.15), y + Inches(0.15), Inches(3.45), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = cat
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_TEAL

        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(15)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_CYAN
        p2.space_before = Pt(2)

        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(12)
        p3.font.color.rgb = COLOR_WHITE
        p3.space_before = Pt(6)

    set_speaker_notes(s14, "We used Python Flask for the backend API, Vanilla JavaScript and CSS for a fast glowing dark-mode frontend, and Chart.js for live analytics.")

    # -------------------------------------------------------------
    # SLIDE 15: Project Workflow (With Threat Nodes Graphic Banner)
    # -------------------------------------------------------------
    s15 = prs.slides.add_slide(blank_layout)
    add_bg(s15)
    add_header(s15, "Project Execution Workflow")
    add_footer(s15, 15)

    wf_steps = [
        ("Step 1: Input Capture", "User types prompt in Prompt Sandbox."),
        ("Step 2: Dual Scanning", "High-speed Regex + Transformer NER spans."),
        ("Step 3: Mask & Score", "Partial Masking (Aadhaar) & Risk score."),
        ("Step 4: Audit & Route", "Sanitized text sent to LLM + SHA-256 Log.")
    ]

    for idx, (w_title, w_desc) in enumerate(wf_steps):
        x = Inches(0.8 + idx * 2.98)
        y = Inches(1.6)

        box = s15.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.75), Inches(2.3))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD
        box.line.color.rgb = COLOR_CYAN
        box.line.width = Pt(1.5)

        tb = s15.shapes.add_textbox(x + Inches(0.1), y + Inches(0.12), Inches(2.55), Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = w_title
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_CYAN
        p1.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = w_desc
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = COLOR_WHITE
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(10)

        if idx < 3:
            arrow = s15.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(2.77), Inches(2.6), Inches(0.2), Inches(0.35))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLOR_CYAN
            arrow.line.fill.background()

    if os.path.exists(IMG_NODES):
        s15.shapes.add_picture(IMG_NODES, Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.5))

    set_speaker_notes(s15, "This flowchart shows the end-to-end user workflow: from prompt submission, dual scanning, partial masking, to final audit logging in real time.")

    # -------------------------------------------------------------
    # SLIDE 16: Advantages of the System
    # -------------------------------------------------------------
    s16 = prs.slides.add_slide(blank_layout)
    add_bg(s16)
    add_header(s16, "Advantages of the System")
    add_footer(s16, 16)

    advs = [
        ("100% Data Privacy Protection:", "Prevents accidental leakage of confidential Indian identity tags and financial data."),
        ("Context-Preserving Masking:", "Keeps AI responses accurate without ruining sentence structure or meaning."),
        ("Ultra-Fast Performance:", "Executes full multi-stage detection in under 50ms with zero noticeable delay."),
        ("Regulatory Compliance Pre-Aligned:", "Supports GDPR, HIPAA, and DPDP Act 2023 data protection standards."),
        ("Interactive User Dashboard:", "Includes live threat charts, Red Team cyber range simulator, and Security Copilot assistant.")
    ]

    for idx, (title, desc) in enumerate(advs):
        box = s16.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6 + idx * 1.02), Inches(11.7), Inches(0.88))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD
        box.line.color.rgb = COLOR_GREEN
        box.line.width = Pt(1)

        tb = s16.shapes.add_textbox(Inches(1.1), Inches(1.68 + idx * 1.02), Inches(11.1), Inches(0.72))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        
        r1 = p.add_run()
        r1.text = "✅ " + title + " "
        r1.font.bold = True
        r1.font.size = Pt(15)
        r1.font.color.rgb = COLOR_GREEN

        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(14)
        r2.font.color.rgb = COLOR_WHITE

    set_speaker_notes(s16, "The main advantage of our system is that it protects privacy without slowing down the AI or ruining the quality of answers. It also helps companies meet government data laws.")

    # -------------------------------------------------------------
    # SLIDE 17: Future Enhancements
    # -------------------------------------------------------------
    s17 = prs.slides.add_slide(blank_layout)
    add_bg(s17)
    add_header(s17, "Future Enhancements")
    add_footer(s17, 17)

    futures = [
        "Multimodal AI Support: Add real-time image OCR and audio scanning for multimodal LLM inputs.",
        "Browser Extension Integration: Develop Chrome / Edge extension to auto-sanitize ChatGPT & Claude web UI.",
        "Advanced Homomorphic Encryption: Expand zero-knowledge encrypted computation for cloud LLM servers.",
        "Hardware Security Module (HSM): Integrate hardware key storage for enterprise-grade server deployment.",
        "Automated Compliance Reports: Add one-click PDF / CSV audit report generator for compliance officers."
    ]

    for idx, fut in enumerate(futures):
        box = s17.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6 + idx * 1.02), Inches(11.7), Inches(0.88))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD
        box.line.color.rgb = COLOR_CYAN
        box.line.width = Pt(1)

        tb = s17.shapes.add_textbox(Inches(1.1), Inches(1.68 + idx * 1.02), Inches(11.1), Inches(0.72))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        parts = fut.split(":", 1)

        r1 = p.add_run()
        r1.text = "🚀 " + parts[0] + ":"
        r1.font.bold = True
        r1.font.size = Pt(15)
        r1.font.color.rgb = COLOR_CYAN

        r2 = p.add_run()
        r2.text = parts[1]
        r2.font.size = Pt(14)
        r2.font.color.rgb = COLOR_WHITE

    set_speaker_notes(s17, "In the future, we plan to add image and voice data scanning, create a browser extension for ChatGPT, and support automated compliance report downloads.")

    # -------------------------------------------------------------
    # SLIDE 18: Conclusion
    # -------------------------------------------------------------
    s18 = prs.slides.add_slide(blank_layout)
    add_bg(s18)
    add_header(s18, "Conclusion")
    add_footer(s18, 18)

    concls = [
        "SecureLLM Shield successfully solves the critical problem of privacy leakage in Large Language Models.",
        "Combines 20-entity PII detection, context-preserving partial masking, and prompt injection defense.",
        "Provides enterprise-grade audit logging with SHA-256 cryptographic chain verification.",
        "Proven 0% bypass rate under automated 1,000-payload Red Team cyber range attack campaigns.",
        "Offers a scalable, ultra-fast security gateway for safe enterprise AI deployment."
    ]

    for idx, conc in enumerate(concls):
        box = s18.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6 + idx * 1.02), Inches(11.7), Inches(0.88))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD
        box.line.color.rgb = COLOR_CYAN
        box.line.width = Pt(1)

        tb = s18.shapes.add_textbox(Inches(1.1), Inches(1.68 + idx * 1.02), Inches(11.1), Inches(0.72))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        
        r1 = p.add_run()
        r1.text = "📌 " + conc
        r1.font.size = Pt(14)
        r1.font.color.rgb = COLOR_WHITE

    set_speaker_notes(s18, "To conclude, SecureLLM Shield provides a complete and reliable security shield for using AI safely in organizations. It bridges the gap between AI convenience and data security.")

    # -------------------------------------------------------------
    # SLIDE 19: Thank You (With Graphic Image)
    # -------------------------------------------------------------
    s19 = prs.slides.add_slide(blank_layout)
    add_bg(s19)

    card19 = s19.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.8), Inches(11.733), Inches(5.9))
    card19.fill.solid()
    card19.fill.fore_color.rgb = COLOR_CARD
    card19.line.color.rgb = COLOR_CYAN
    card19.line.width = Pt(2)

    tb19 = s19.shapes.add_textbox(Inches(1.1), Inches(1.4), Inches(7.2), Inches(4.8))
    tf19 = tb19.text_frame
    tf19.word_wrap = True

    p = tf19.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    p2 = tf19.add_paragraph()
    p2.text = "Questions & Discussion"
    p2.font.size = Pt(22)
    p2.font.color.rgb = COLOR_WHITE
    p2.space_before = Pt(10)

    p3 = tf19.add_paragraph()
    p3.text = "SecureLLM Shield: Privacy Threat Detection & Migration Framework for LLMs"
    p3.font.size = Pt(13)
    p3.font.color.rgb = COLOR_TEAL
    p3.space_before = Pt(20)

    p4 = tf19.add_paragraph()
    p4.text = "Project Guide: Mr. G. Saravanan Gokul, M.E. (Assistant Professor, CSE-AIML)\nTeam Members: BARATHKUMAR M & SUBASH P\nS.A. Engineering College (Autonomous)"
    p4.font.size = Pt(13)
    p4.font.bold = True
    p4.font.color.rgb = COLOR_WHITE
    p4.space_before = Pt(12)

    p5 = tf19.add_paragraph()
    p5.text = "GitHub Repository: https://github.com/subashp-2104/SecureLLM-Shield"
    p5.font.size = Pt(11.5)
    p5.font.color.rgb = COLOR_MUTED
    p5.space_before = Pt(14)

    if os.path.exists(IMG_SHIELD):
        s19.shapes.add_picture(IMG_SHIELD, Inches(8.5), Inches(1.3), Inches(3.7), Inches(4.9))

    set_speaker_notes(s19, "Thank you respected evaluators for your time and guidance. We are now open for any questions and feedback.")

    output_path = os.path.join(base_dir, "SecureLLM_Shield_Presentation.pptx")
    try:
        prs.save(output_path)
        print(f"SUCCESS: Rebuilt 19-slide presentation with repaired Slide 7 at {output_path}")
    except PermissionError:
        output_alt = os.path.join(base_dir, "SecureLLM_Shield_Presentation_Updated.pptx")
        prs.save(output_alt)
        print(f"SUCCESS: Rebuilt 19-slide presentation with repaired Slide 7 at {output_alt}")

if __name__ == "__main__":
    build_presentation()
